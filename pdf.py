from flask import Flask, request, send_file
import os
import pdfkit
from werkzeug.utils import secure_filename
from docx import Document
import pandas as pd
from pptx import Presentation

app = Flask(__name__)
UPLOAD_FOLDER = "uploads"
OUTPUT_FOLDER = "output"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

def convert_docx_to_pdf(file_path, output_path):
    doc = Document(file_path)
    html_content = "<html><body>"
    for para in doc.paragraphs:
        html_content += f"<p>{para.text}</p>"
    html_content += "</body></html>"
    with open("temp.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    pdfkit.from_file("temp.html", output_path)

def convert_excel_to_pdf(file_path, output_path):
    df = pd.read_excel(file_path)
    html_content = df.to_html()
    with open("temp.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    pdfkit.from_file("temp.html", output_path)

def convert_pptx_to_pdf(file_path, output_path):
    prs = Presentation(file_path)
    html_content = "<html><body>"
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                html_content += f"<p>{shape.text}</p>"
    html_content += "</body></html>"
    with open("temp.html", "w", encoding="utf-8") as f:
        f.write(html_content)
    pdfkit.from_file("temp.html", output_path)

@app.route("/convert", methods=["POST"])
def convert_file():
    if 'file' not in request.files:
        return {"error": "No file provided"}, 400
    
    file = request.files['file']
    filename = secure_filename(file.filename)
    file_ext = filename.split(".")[-1].lower()
    file_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(file_path)
    
    output_filename = f"{filename.rsplit('.', 1)[0]}.pdf"
    output_path = os.path.join(OUTPUT_FOLDER, output_filename)
    
    if file_ext == "docx":
        convert_docx_to_pdf(file_path, output_path)
    elif file_ext in ["xls", "xlsx"]:
        convert_excel_to_pdf(file_path, output_path)
    elif file_ext == "pptx":
        convert_pptx_to_pdf(file_path, output_path)
    else:
        return {"error": "Unsupported file format"}, 400
    
    return send_file(output_path, as_attachment=True)

if __name__ == "__main__":
    app.run(debug=True)
